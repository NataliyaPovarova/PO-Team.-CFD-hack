import pandas as pd
import os
from datetime import date, datetime
import matplotlib
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import sklearn
from functools import reduce
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from pptx import Presentation


def clustering(data):
    '''takes xlsx file as an input,
    saves png with colored clusters'''
    data = data.copy()
    del data['Unnamed: 0']
    data.drop(['ФИО', 'Роль в мероприятии'], axis=1, inplace=True)
    text_columns = ['Список компетенций', 'Должность', 'Категория', 'Место работы', 'Профессия', 'Образование',
                'Место образования', 'Специальность']
    data['Трудовой стаж'] = current_year - pd.DatetimeIndex(data['Начало трудового стажа']).year
    data['Возраст'] = current_year - pd.DatetimeIndex(data['Дата рождения']).year
    data['Стаж работы в РОСАТОМ'] = current_year - pd.DatetimeIndex(data['Начало трудовой деятельности в РОСАТОМ']).year
    data.drop(['Начало трудового стажа', 'Дата рождения', 'Начало трудовой деятельности в РОСАТОМ', 'Год оканчания'],
          axis=1, inplace=True)
    data.fillna({'Должность':'нет данных', 'Категория':'нет данных', 'Место работы': 'нет данных',
             'Профессия': 'нет данных', 'Образование': 'нет данных', 'Место образования': 'нет данных',
             'Специальность': 'нет данных'},
            inplace=True)
    data.fillna({'Трудовой стаж': data['Трудовой стаж'].median(),
            'Возраст': data['Возраст'].median(),
            'Стаж работы в РОСАТОМ': data['Стаж работы в РОСАТОМ'].median()},
           inplace=True)
    data.fillna(0, inplace=True)
    encoded_data = pd.get_dummies(data=data, columns=text_columns)
    norm_data = (encoded_data - encoded_data.min())/(encoded_data.max() - encoded_data.min())
    kmeans_model = KMeans(n_clusters = 3)
    kmeans_model.fit(norm_data)
    encoded_data["clusters"] = kmeans_model.labels_
    y = encoded_data['clusters']
    X = encoded_data.loc[:, encoded_data.columns != 'clusters']
    X_norm = (X - X.min())/(X.max() - X.min())
    pca = PCA(n_components=2) #2-dimensional PCA
    transformed = pd.DataFrame(pca.fit_transform(X_norm))
    plt.scatter(transformed[y==0][0], transformed[y==0][1], label='Cluster 1', c='red')
    plt.scatter(transformed[y==1][0], transformed[y==1][1], label='Cluster 2', c='blue')
    plt.scatter(transformed[y==2][0], transformed[y==2][1], label='Cluster 3', c='lightgreen')

    plt.legend()
    plt.savefig('АП\\clusters.png')


def calculating_age(birthdate):
    if birthdate == None:
        pass
    else:
        today = date.today()
        age = today.year - birthdate.year - ((today.month, today.day) < (birthdate.month, birthdate.day))
        return age


def contest_distribution(df):
    grades_starting_column = df.columns.get_loc('Аддитивные технологии')
    grades_final_column = df.columns.get_loc('Электроника')
    competition_df = df.iloc[:, grades_starting_column:grades_final_column + 1]
    counted_competitions_df = competition_df.count()
    counted_competitions_df = counted_competitions_df.nlargest(10)

    counted_competitions_plot = counted_competitions_df.plot.bar(color=colour_pallette,
                                                                 ylabel='Количество участников')
    counted_competitions_fig = counted_competitions_plot.get_figure()
    counted_competitions_fig.savefig('АП\\counted_competitions.png', bbox_inches='tight')


def mean_contest(df):
    grades_starting_column = df.columns.get_loc('Аддитивные технологии')
    grades_final_column = df.columns.get_loc('Электроника')
    competition_df = df.iloc[:, grades_starting_column:grades_final_column + 1]
    counted_competitions_df = competition_df.count()
    mean_competition_df = competition_df.mean(axis=0)
    mean_competition_df = mean_competition_df.nlargest(10)

    mean_competition_plot = mean_competition_df.plot.bar(color=colour_pallette,
                                                         ylabel='Средний балл')
    mean_competition_fig = mean_competition_plot.get_figure()
    mean_competition_fig.savefig('АП\\mean_competition.png', bbox_inches='tight')


def age_distribution(df):
    unclean_birthdates_df = df['Дата рождения']
    birthdates_df = unclean_birthdates_df.dropna()
    age_df = birthdates_df.map(calculating_age)
    first_age_group = len(age_df[(age_df > 0) & (age_df <= 21)])
    second_age_group = len(age_df[(age_df > 21) & (age_df <= 28)])
    third_age_group = len(age_df[(age_df > 28) & (age_df <= 35)])
    fourth_age_group = len(age_df[(age_df > 35) & (age_df <= 45)])
    fifth_age_group = len(age_df[(age_df > 45) & (age_df <= 55)])
    sixth_age_group = len(age_df[age_df > 55])
    age_groups_data = [first_age_group, second_age_group, third_age_group, fourth_age_group, fifth_age_group,
                       sixth_age_group]
    age_groups_df = pd.DataFrame(age_groups_data, index=['17-21', '22-28', '29-35', '36-45', '46-55', 'выше 55'],
                                 columns=['Возрастные группы'])

    age_group_plot = age_groups_df.plot.pie(y='Возрастные группы', figsize=(16, 9), fontsize=14, legend=False)
    age_group_fig = age_group_plot.get_figure()
    age_group_fig.savefig('АП\\age_groups.png')


def rosatom_experience(df):
    today = date.today()
    unclean_experience_df = df['Начало трудовой деятельности в РОСАТОМ']
    rosatom_experience_df = unclean_experience_df.dropna()
    rosatom_start_year_df = rosatom_experience_df.dt.year
    today = date.today()
    experience_difference_df = today.year - rosatom_start_year_df

    first_experience_group = len(experience_difference_df[experience_difference_df < 3])
    second_experience_group = len(
        experience_difference_df[(experience_difference_df >= 3) & (experience_difference_df <= 5)])
    third_experience_group = len(
        experience_difference_df[(experience_difference_df > 5) & (experience_difference_df <= 10)])
    fourth_experience_group = len(
        experience_difference_df[(experience_difference_df > 10) & (experience_difference_df <= 15)])
    fifth_experience_group = len(
        experience_difference_df[(experience_difference_df > 15) & (experience_difference_df <= 20)])
    sixth_experience_group = len(experience_difference_df[experience_difference_df > 20])
    experience_groups_data = [first_experience_group, second_experience_group, third_experience_group,
                              fourth_experience_group, fifth_experience_group, sixth_experience_group]
    experience_groups_df = pd.DataFrame(experience_groups_data,
                                        index=['менее 3 лет', '3-5 лет', '6-10 лет', '11-15 лет',
                                               '15-20 лет', 'выше 20 лет'], columns=['Опыт работы'])

    experience_group_plot = experience_groups_df.plot.pie(y='Опыт работы', figsize=(16, 9), fontsize=14, legend=False)
    experience_group_fig = experience_group_plot.get_figure()
    experience_group_fig.savefig('АП\\experience_groups.png')


def gender_skills(df):
    '''Пол с компетенциями/соревнованиями
    takes dataframe as input, saves heatmap as png'''
    counted_gender_comp = df[['Пол', 'Список компетенций']].value_counts().reset_index(name='count')
    heatmap = counted_gender_comp.pivot_table(index='Пол',
            columns='Список компетенций',
            values='count')
    heatmap = heatmap.T
    heatmap.fillna(0, inplace=True)
    heatmap = heatmap.astype('int64')
    heatmap.reset_index(inplace=True)
    heatmap['Список компетенций'] = heatmap['Список компетенций'].str.replace(';', '')
    heatmap['sum'] = heatmap[0] + heatmap[1]
    heatmap.sort_values(by='sum', ascending=False, inplace=True)
    heatmap.set_index('Список компетенций', inplace=True)
    del heatmap['sum']
    top_heatmap = heatmap.iloc[:10, :]

    heatmap_5 = sns.heatmap(top_heatmap, annot=True, cmap="YlGnBu")
    heatmap_5_fig = heatmap_5.get_figure()
    heatmap_5_fig.savefig('АП\\Пол с компетенциями.png', dpi=300, bbox_inches='tight')


def gender_rosatom_start(df):
    '''Пол с началом деятельности в Росатоме
    takes dataframe as input, saves heatmap as png'''
    rosatom_df = df[~df['Начало трудовой деятельности в РОСАТОМ'].isnull()]
    rosatom_df.reset_index(inplace=True)
    rosatom_df['Год начала работы в РОСАТОМ'] = pd.DatetimeIndex(rosatom_df['Начало трудовой деятельности в РОСАТОМ']).year
    rosatom_dynamics = rosatom_df[['Год начала работы в РОСАТОМ', 'Пол']].value_counts().reset_index(name='count')
    rosatom_heatmap = rosatom_dynamics.pivot_table(index='Год начала работы в РОСАТОМ',
            columns='Пол',
            values='count')
    rosatom_heatmap.reset_index(inplace=True)
    rosatom_heatmap.sort_values(by='Год начала работы в РОСАТОМ', ascending=True, inplace=True)
    rosatom_heatmap.set_index('Год начала работы в РОСАТОМ', inplace=True)
    rosatom_heatmap.fillna(0, inplace=True)
    heatmap_6 = rosatom_heatmap.plot.line()
    heatmap_6_fig = heatmap_6.get_figure()
    heatmap_6_fig.savefig('АП\\Пол с началом деятельности в Росатоме.png', dpi=300, bbox_inches='tight')


def ed_skills(df):
    '''Компетенции с образованием (тепловая карта)
    takes a df, saves a heatmap'''
    skills_education = df[['Список компетенций', 'Образование']].value_counts().reset_index(name='count')
    ed_heatmap = skills_education.pivot_table(index='Образование',
            columns='Список компетенций',
            values='count')
    ed_heatmap = ed_heatmap.T
    ed_heatmap.fillna(0, inplace=True)
    ed_heatmap.reset_index(inplace=True)
    ed_heatmap['Список компетенций'] = ed_heatmap['Список компетенций'].str.replace(';', '')
    ed_heatmap['sum'] = ed_heatmap['высшее'] + ed_heatmap['общее'] + ed_heatmap['среднее']
    ed_heatmap.sort_values(by='sum', ascending=False)
    ed_heatmap.set_index('Список компетенций', inplace=True)
    top_ed_heatmap = ed_heatmap.iloc[:20, :]
    heatmap_7 = sns.heatmap(top_ed_heatmap, annot=True, cmap="YlGnBu")
    heatmap_7_fig = heatmap_7.get_figure()
    heatmap_7_fig.savefig('АП\\Компетенции с образованием (тепловая карта).png', dpi=300, bbox_inches='tight')



#constants
colour_pallette = ['#000000', '#003274', '#191919', '#025EA1', '#54585A', '#4596D1']
current_date = datetime.today()
current_year = current_date.year

#start of program (applying functions)
period = 10 # last 10 years for dynamics
file_name = 'result.xlsx'
main_df = pd.read_excel(file_name)

get_clustering = clustering(main_df)
get_contest_distribution = contest_distribution(main_df)
get_mean_contest = mean_contest(main_df)
get_age_distribution = age_distribution(main_df)
get_rosatom_experience = rosatom_experience(main_df)
get_gender_skills = gender_skills(main_df)
get_gender_rosatom_start = gender_rosatom_start(main_df)
get_ed_skills = ed_skills(main_df)


#at the end - inserting resulted images at our automated powerpoint dashboard
prs = Presentation('PO_Team.Rosatom.pptx')

##this dict and number of charts could change depending on actual needs/datasets
analytics_dict = {'Кластеризация' : 'clusters.png',
                  'Статистика с баллами по компетенциям' : 'counted_competitions.png',
                  'Средний балл по компетенциям' : 'mean_competition.png',
                  'Представленные демографические группы' : 'age_groups.png',
                  'Трудовой стаж в Росатоме' : 'experience_groups.png',
                  'Соотношение гендера с компетенциями/соревнованиями' : 'Пол с компетенциями.png',
                  'Соотношение гендера с началом работы в Росатоме' : 'Пол с началом деятельности в Росатоме.png',
                  'Соотношение компетенций с указанным образованием' : 'Компетенции с образованием (тепловая карта).png'}

for title, chart_name in analytics_dict.items():
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    text_placeholder = slide.placeholders[1]
    text_placeholder.text = title

    picture_placeholder = slide.placeholders[13]
    picture = picture_placeholder.insert_picture(os.path.join('АП', chart_name))

new_prs = prs.save('PO_Team.Rosatom_modified.pptx')

print('finale')

